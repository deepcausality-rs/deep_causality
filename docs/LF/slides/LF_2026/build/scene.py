"""A measured scene graph with two backends: PowerPoint and PNG preview.

The deck is described once as primitives. `to_pptx` writes the .pptx; `to_png`
rasterises the same primitives with PIL so the preview cannot drift from the
deliverable.

Text is measured against the real TTFs and pre-wrapped, so a block's height is
`len(lines) * leading` exactly — no autofit, no guessing, no overlap.
"""

import os
from PIL import Image, ImageDraw, ImageFont

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
FONT_DIR = os.path.expanduser("~/Library/Fonts")

FONT_FILES = {
    "Geist": "Geist-Regular.ttf",
    "Geist Light": "Geist-Light.ttf",
    "Geist Medium": "Geist-Medium.ttf",
    "Geist SemiBold": "Geist-SemiBold.ttf",
    "Geist Bold": "Geist-Bold.ttf",
    "JetBrains Mono": "JetBrainsMono-Regular.ttf",
    "JetBrains Mono Medium": "JetBrainsMono-Medium.ttf",
    "JetBrains Mono Bold": "JetBrainsMono-Bold.ttf",
}

_MEAS_SCALE = 16  # measure at 16 px/pt for sub-point accuracy
_cache = {}


def _font(name, px):
    key = (name, int(px))
    if key not in _cache:
        _cache[key] = ImageFont.truetype(
            os.path.join(FONT_DIR, FONT_FILES[name]), int(px))
    return _cache[key]


def text_pt(s, name, size, spc=0):
    """Width of `s` in points. `spc` is DrawingML tracking, hundredths of a pt."""
    if not s:
        return 0.0
    f = _font(name, size * _MEAS_SCALE)
    return f.getlength(s) / _MEAS_SCALE + (spc / 100.0) * len(s)


def text_in(s, name, size, spc=0):
    return text_pt(s, name, size, spc) / 72.0


def wrap(text, name, size, max_w_in, spc=0, safety=0.985):
    """Pre-wrap to a hard line list. Honours explicit newlines."""
    limit = max_w_in * safety
    out = []
    for chunk in str(text).split("\n"):
        words = chunk.split()
        if not words:
            out.append("")
            continue
        cur = words[0]
        for w in words[1:]:
            trial = cur + " " + w
            if text_in(trial, name, size, spc) <= limit:
                cur = trial
            else:
                out.append(cur)
                cur = w
        out.append(cur)
    return out


# --------------------------------------------------------------------------
class Scene:
    """One slide's primitives, in paint order."""

    def __init__(self, bg):
        self.bg = bg
        self.ops = []
        self.notes = ""

    # -- primitives -------------------------------------------------------
    def rect(self, x, y, w, h, fill=None, fill_alpha=None, line=None,
             line_w=0.75, radius=0.0, gradient=None):
        self.ops.append(("rect", dict(x=x, y=y, w=w, h=h, fill=fill,
                                      fill_alpha=fill_alpha, line=line,
                                      line_w=line_w, radius=radius,
                                      gradient=gradient)))

    def diamond(self, x, y, w, h, fill=None, line=None, line_w=0.9):
        self.ops.append(("diamond", dict(x=x, y=y, w=w, h=h, fill=fill,
                                         line=line, line_w=line_w)))

    def seg(self, x1, y1, x2, y2, color, w=0.75, alpha=None, arrow=False):
        self.ops.append(("seg", dict(x1=x1, y1=y1, x2=x2, y2=y2, color=color,
                                     w=w, alpha=alpha, arrow=arrow)))

    def oval(self, cx, cy, d, fill, alpha=None):
        self.ops.append(("oval", dict(cx=cx, cy=cy, d=d, fill=fill,
                                      alpha=alpha)))

    def image(self, path, x, y, w, h):
        self.ops.append(("image", dict(path=path, x=x, y=y, w=w, h=h)))

    def lines(self, x, y, w, lines, name, size, color, leading,
              align="l", spc=0, runs=None):
        """A pre-wrapped text block. Height is exactly len(lines)*leading."""
        self.ops.append(("text", dict(x=x, y=y, w=w, lines=lines, name=name,
                                      size=size, color=color,
                                      leading=leading, align=align, spc=spc,
                                      runs=runs)))


# --------------------------------------------------------------------------
# PPTX backend
# --------------------------------------------------------------------------
def _rgb(h):
    return RGBColor(h[0], h[1], h[2])


def _alpha(shape, pct, which="fill"):
    spPr = shape._element.spPr
    node = spPr.find(qn("a:solidFill")) if which == "fill" else spPr.find(qn("a:ln"))
    if node is None:
        return
    if which == "line":
        node = node.find(qn("a:solidFill"))
        if node is None:
            return
    clr = node.find(qn("a:srgbClr"))
    if clr is None:
        return
    for old in clr.findall(qn("a:alpha")):
        clr.remove(old)
    clr.append(parse_xml(f'<a:alpha xmlns:a="{NS_A}" val="{int(pct * 1000)}"/>'))


_ALIGN = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}


def scene_to_slide(scene, prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(scene.bg)

    for kind, o in scene.ops:
        if kind == "rect":
            shp = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE if o["radius"] else MSO_SHAPE.RECTANGLE,
                Inches(o["x"]), Inches(o["y"]), Inches(o["w"]), Inches(o["h"]))
            if o["radius"]:
                shp.adjustments[0] = min(0.5, o["radius"] / min(o["w"], o["h"]))
            if o["gradient"]:
                f = shp.fill
                f.gradient()
                f.gradient_stops[0].color.rgb = _rgb(o["gradient"][0])
                f.gradient_stops[0].position = 0.0
                f.gradient_stops[1].color.rgb = _rgb(o["gradient"][1])
                f.gradient_stops[1].position = 1.0
                f.gradient_angle = 90.0
            elif o["fill"]:
                shp.fill.solid()
                shp.fill.fore_color.rgb = _rgb(o["fill"])
                if o["fill_alpha"] is not None:
                    _alpha(shp, o["fill_alpha"], "fill")
            else:
                shp.fill.background()
            if o["line"]:
                shp.line.color.rgb = _rgb(o["line"])
                shp.line.width = Pt(o["line_w"])
            else:
                shp.line.fill.background()
            shp.shadow.inherit = False

        elif kind == "diamond":
            shp = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(o["x"]),
                                         Inches(o["y"]), Inches(o["w"]),
                                         Inches(o["h"]))
            if o["fill"]:
                shp.fill.solid()
                shp.fill.fore_color.rgb = _rgb(o["fill"])
            else:
                shp.fill.background()
            if o["line"]:
                shp.line.color.rgb = _rgb(o["line"])
                shp.line.width = Pt(o["line_w"])
            else:
                shp.line.fill.background()
            shp.shadow.inherit = False

        elif kind == "seg":
            c = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, Inches(o["x1"]), Inches(o["y1"]),
                Inches(o["x2"]), Inches(o["y2"]))
            c.line.color.rgb = _rgb(o["color"])
            c.line.width = Pt(o["w"])
            if o["alpha"] is not None:
                _alpha(c, o["alpha"], "line")
            if o["arrow"]:
                c._element.spPr.find(qn("a:ln")).append(parse_xml(
                    f'<a:tailEnd xmlns:a="{NS_A}" type="triangle" w="med" len="med"/>'))

        elif kind == "oval":
            d = o["d"]
            shp = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(o["cx"] - d / 2), Inches(o["cy"] - d / 2),
                Inches(d), Inches(d))
            shp.fill.solid()
            shp.fill.fore_color.rgb = _rgb(o["fill"])
            if o["alpha"] is not None:
                _alpha(shp, o["alpha"], "fill")
            shp.line.fill.background()
            shp.shadow.inherit = False

        elif kind == "image":
            slide.shapes.add_picture(o["path"], Inches(o["x"]), Inches(o["y"]),
                                     Inches(o["w"]), Inches(o["h"]))

        elif kind == "text":
            n = max(1, len(o["lines"]))
            tb = slide.shapes.add_textbox(
                Inches(o["x"]), Inches(o["y"]), Inches(o["w"]),
                Inches(n * o["leading"] / 72.0 + 0.02))
            tf = tb.text_frame
            tf.word_wrap = False
            tf.auto_size = MSO_AUTO_SIZE.NONE
            tf.margin_left = tf.margin_right = 0
            tf.margin_top = tf.margin_bottom = 0
            tf.vertical_anchor = MSO_ANCHOR.TOP
            runs_per_line = o["runs"] or [None] * n
            for i, ln in enumerate(o["lines"]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.alignment = _ALIGN[o["align"]]
                p.line_spacing = Pt(o["leading"])
                p.space_before = Pt(0)
                p.space_after = Pt(0)
                spec = runs_per_line[i] if i < len(runs_per_line) else None
                pieces = spec if spec else [(ln, o["name"], o["color"])]
                for txt, fname, col in pieces:
                    r = p.add_run()
                    r.text = txt
                    r.font.name = fname
                    r.font.size = Pt(o["size"])
                    r.font.color.rgb = _rgb(col)
                    if o["spc"]:
                        r.font._rPr.set("spc", str(int(o["spc"])))

    if scene.notes:
        slide.notes_slide.notes_text_frame.text = scene.notes
    return slide


# --------------------------------------------------------------------------
# PNG backend
# --------------------------------------------------------------------------
def scene_to_png(scene, path, w_in, h_in, ppi=120):
    W, H = int(w_in * ppi), int(h_in * ppi)
    img = Image.new("RGB", (W, H), scene.bg)

    def I(v):
        return v * ppi

    def blend(base, color, alpha_pct):
        a = alpha_pct / 100.0
        return tuple(int(base[i] * (1 - a) + color[i] * a) for i in range(3))

    for kind, o in scene.ops:
        ov = Image.new("RGBA", (W, H), (0, 0, 0, 0))
        d = ImageDraw.Draw(ov)
        if kind == "rect":
            box = [I(o["x"]), I(o["y"]), I(o["x"] + o["w"]), I(o["y"] + o["h"])]
            fill = None
            if o["gradient"]:
                fill = None  # painted below as a vertical ramp
            elif o["fill"]:
                a = 255 if o["fill_alpha"] is None else int(o["fill_alpha"] * 2.55)
                fill = o["fill"] + (a,)
            outline = o["line"] + (255,) if o["line"] else None
            lw = max(1, int(round(o["line_w"] * ppi / 72.0)))
            r = int(I(o["radius"]))
            if o["gradient"]:
                g = Image.new("RGB", (1, max(2, int(I(o["h"])))))
                gd = ImageDraw.Draw(g)
                c0, c1 = o["gradient"]
                hh = g.size[1]
                for yy in range(hh):
                    t = yy / max(1, hh - 1)
                    gd.point((0, yy), tuple(int(c0[i] + (c1[i] - c0[i]) * t)
                                            for i in range(3)))
                g = g.resize((max(1, int(I(o["w"]))), max(2, int(I(o["h"])))))
                mask = Image.new("L", g.size, 0)
                ImageDraw.Draw(mask).rounded_rectangle(
                    [0, 0, g.size[0] - 1, g.size[1] - 1], radius=r, fill=255)
                img.paste(g, (int(I(o["x"])), int(I(o["y"]))), mask)
            if r:
                d.rounded_rectangle(box, radius=r, fill=fill, outline=outline,
                                    width=lw)
            else:
                d.rectangle(box, fill=fill, outline=outline, width=lw)

        elif kind == "diamond":
            cx, cy = I(o["x"] + o["w"] / 2), I(o["y"] + o["h"] / 2)
            pts = [(cx, I(o["y"])), (I(o["x"] + o["w"]), cy),
                   (cx, I(o["y"] + o["h"])), (I(o["x"]), cy)]
            d.polygon(pts, fill=o["fill"] + (255,) if o["fill"] else None,
                      outline=o["line"] + (255,) if o["line"] else None,
                      width=max(1, int(round(o["line_w"] * ppi / 72.0))))

        elif kind == "seg":
            col = o["color"] + (int((o["alpha"] or 100) * 2.55),)
            lw = max(1, int(round(o["w"] * ppi / 72.0)))
            d.line([I(o["x1"]), I(o["y1"]), I(o["x2"]), I(o["y2"])], fill=col,
                   width=lw)
            if o["arrow"]:
                import math
                ang = math.atan2(o["y2"] - o["y1"], o["x2"] - o["x1"])
                L, Wd = I(0.10), I(0.052)
                tipx, tipy = I(o["x2"]), I(o["y2"])
                bx, by = tipx - L * math.cos(ang), tipy - L * math.sin(ang)
                d.polygon([(tipx, tipy),
                           (bx - Wd * math.sin(ang), by + Wd * math.cos(ang)),
                           (bx + Wd * math.sin(ang), by - Wd * math.cos(ang))],
                          fill=col)

        elif kind == "oval":
            r = I(o["d"]) / 2
            d.ellipse([I(o["cx"]) - r, I(o["cy"]) - r, I(o["cx"]) + r,
                       I(o["cy"]) + r],
                      fill=o["fill"] + (int((o["alpha"] or 100) * 2.55),))

        elif kind == "image":
            im = Image.open(o["path"]).convert("RGBA")
            im = im.resize((max(1, int(I(o["w"]))), max(1, int(I(o["h"])))),
                           Image.LANCZOS)
            ov.alpha_composite(im, (int(I(o["x"])), int(I(o["y"]))))

        elif kind == "text":
            f = _font(o["name"], o["size"] * ppi / 72.0)
            asc, _ = f.getmetrics()
            lead_px = o["leading"] * ppi / 72.0
            runs_per_line = o["runs"] or [None] * len(o["lines"])
            for i, ln in enumerate(o["lines"]):
                spec = runs_per_line[i] if i < len(runs_per_line) else None
                pieces = spec if spec else [(ln, o["name"], o["color"])]
                # Measure with the very advances the draw loop will use, so a
                # tracked run centres on screen exactly where it centres in the
                # file. PIL hints per pixel size; the measuring size differs.
                adv_px = (o["spc"] / 100.0) * ppi / 72.0
                total = 0.0
                for t, fn, _ in pieces:
                    ff = _font(fn, o["size"] * ppi / 72.0)
                    total += (sum(ff.getlength(c) for c in t) + adv_px * len(t)
                              if o["spc"] else ff.getlength(t))
                total -= adv_px if o["spc"] else 0
                if o["align"] == "c":
                    px = I(o["x"]) + (I(o["w"]) - total) / 2
                elif o["align"] == "r":
                    px = I(o["x"]) + I(o["w"]) - total
                else:
                    px = I(o["x"])
                py = I(o["y"]) + i * lead_px + (lead_px - asc * 1.0) * 0.55
                for txt, fname, col in pieces:
                    ff = _font(fname, o["size"] * ppi / 72.0)
                    if o["spc"]:
                        adv = (o["spc"] / 100.0) * ppi / 72.0
                        for ch in txt:
                            d.text((px, py), ch, font=ff, fill=col + (255,))
                            px += ff.getlength(ch) + adv
                    else:
                        d.text((px, py), txt, font=ff, fill=col + (255,))
                        px += ff.getlength(txt)

        img = Image.alpha_composite(img.convert("RGBA"), ov).convert("RGB")

    img.save(path)
    return path
